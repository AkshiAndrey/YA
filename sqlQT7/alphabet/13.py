from pptx import Presentation
# создаем новую презентацию
prs = Presentation()
# получаем схему расположения элементов для заголовочного слайда
title_slide_layout = prs.slide_layouts[0]
# создаем заголовочный слайд
slide = prs.slides.add_slide(title_slide_layout)
# создаем у слайда заголовок и текст
title = slide.shapes.title
print("Контейнер для текста:", slide.placeholders[1])
subtitle = slide.placeholders[1]

title.text = "Тестовый заголовок"
subtitle.text = "Тестовый текст"
# создаем новый слайд со схемой для добавления изображений
slide = prs.slides.add_slide(prs.slide_layouts[8])
slide.shapes.title.text = "А теперь с картинкой"
# добавляем изображение
print("Контейнер для изображения:", slide.placeholders[1])
placeholder = slide.placeholders[1]
placeholder.insert_picture('1jpg.jpg')
# сохраняем презентацию
prs.save('test.pptx')